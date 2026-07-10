"""Text extraction from PDF, DOCX, and PPTX files.

Parser libraries ship with the base ``grinta`` install (``pypdf``, ``python-docx``,
``python-pptx``, ``pylatexenc``).
"""

from __future__ import annotations

_MISSING_PARSER_HINT = (
    'Document parser dependency missing. '
    'Reinstall grinta or run: pip install pypdf python-docx python-pptx pylatexenc'
)


def _read_pdf_reader():
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError(_MISSING_PARSER_HINT) from exc
    return PdfReader


def extract_pdf_text(file_path: str) -> str:
    """Return plain text extracted from a PDF file."""
    reader = _read_pdf_reader()(file_path)
    parts: list[str] = []
    for page_idx, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ''
        if text.strip():
            parts.append(f'@@ Page {page_idx} @@\n{text}')
    return '\n\n'.join(parts) if parts else '(empty PDF)'


def extract_docx_text(file_path: str) -> str:
    """Return plain text extracted from a DOCX file."""
    try:
        import docx  # type: ignore[import-untyped, import-not-found]
    except ImportError as exc:
        raise RuntimeError(_MISSING_PARSER_HINT) from exc
    document = docx.Document(file_path)
    parts: list[str] = []
    for para in document.paragraphs:
        text = para.text.strip()
        if text:
            parts.append(text)
    return '\n'.join(parts) if parts else '(empty document)'


def extract_pptx_text(file_path: str) -> str:
    """Return plain text extracted from a PPTX file."""
    try:
        from pptx import Presentation  # type: ignore[import-untyped, import-not-found]
    except ImportError as exc:
        raise RuntimeError(_MISSING_PARSER_HINT) from exc
    presentation = Presentation(file_path)
    parts: list[str] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        slide_lines: list[str] = [f'@@ Slide {slide_idx} @@']
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text:
                slide_lines.append(shape.text)
        if len(slide_lines) > 1:
            parts.append('\n'.join(slide_lines))
    return '\n\n'.join(parts) if parts else '(empty presentation)'
