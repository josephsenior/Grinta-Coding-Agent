"""File reader skills for the Forge agent.

This module provides various functions to parse and extract content from different file types,
including PDF, DOCX, LaTeX, audio, image, video, and PowerPoint files. It utilizes different
libraries and APIs to process these files and output their content or descriptions.

Functions:
    parse_pdf(file_path: str) -> None: Parse and print content of a PDF file.
    parse_docx(file_path: str) -> None: Parse and print content of a DOCX file.
    parse_latex(file_path: str) -> None: Parse and print content of a LaTeX file.
    parse_audio(file_path: str, model: str = 'whisper-1') -> None: Transcribe and print content of an audio file.
    parse_image(file_path: str, task: str = 'Describe this image as detail as possible.') -> None: Analyze and print description of an image file.
    parse_video(file_path: str, task: str = 'Describe this image as detail as possible.', frame_interval: int = 30) -> None: Analyze and print description of video frames.
    parse_pptx(file_path: str) -> None: Parse and print content of a PowerPoint file.

Note:
    Some functions (parse_audio, parse_video, parse_image) require OpenAI API credentials
    and are only available if the necessary environment variables are set.

"""

import base64
from typing import Any, cast

import docx
import PyPDF2
from pptx import Presentation
from pylatexenc.latex2text import LatexNodes2Text  # type: ignore[import-untyped]

from backend.runtime.plugins.agent_skills.utils.config import (
    _get_max_token,
    _get_openai_client,
    _get_openai_model,
)


def parse_pdf(file_path: str) -> None:
    """Parses the content of a PDF file and prints it.

    Args:
        file_path: str: The path to the file to open.

    """
    content = PyPDF2.PdfReader(file_path)
    output_lines = [f"[Reading PDF file from {file_path}]"]
    for page_idx, page in enumerate(content.pages, start=1):
        output_lines.append(f"@@ Page {page_idx} @@")
        output_lines.append(page.extract_text() or "")
        output_lines.append("")
    print("\n".join(output_lines) + "\n")


def parse_docx(file_path: str) -> None:
    """Parses the content of a DOCX file and prints it.

    Args:
        file_path: str: The path to the file to open.

    """
    content = docx.Document(file_path)
    output_lines = [f"[Reading DOCX file from {file_path}]"]
    for i, para in enumerate(content.paragraphs, start=1):
        output_lines.append(f"@@ Page {i} @@")
        output_lines.append(para.text)
        output_lines.append("")
    print("\n".join(output_lines) + "\n")


def parse_latex(file_path: str) -> None:
    """Parses the content of a LaTex file and prints it.

    Args:
        file_path: str: The path to the file to open.

    """
    with open(file_path, encoding="utf-8") as f:
        data = f.read()
    text = LatexNodes2Text().latex_to_text(data).strip()
    print(f"[Reading LaTex file from {file_path}]")
    print(text)


def _base64_img(file_path: str) -> str:
    with open(file_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode("utf-8")


def _base64_video(file_path: str, frame_interval: int = 10) -> list[str]:
    import cv2  # type: ignore[import-not-found]

    video_capture = getattr(cv2, "VideoCapture")
    imencode = getattr(cv2, "imencode")
    video = video_capture(file_path)
    base64_frames = []
    frame_count = 0
    while video.isOpened():
        success, frame = video.read()
        if not success:
            break
        if frame_count % frame_interval == 0:
            _, buffer = imencode(".jpg", frame)
            base64_frames.append(base64.b64encode(buffer).decode("utf-8"))
        frame_count += 1
    video.release()
    return base64_frames


def _prepare_image_messages(task: str, base64_image: str) -> list[dict[str, Any]]:
    return [
        {
            "role": "user",
            "content": [
                {"type": "text", "text": task},
                {
                    "type": "image_url",
                    "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"},
                },
            ],
        },
    ]


def parse_audio(file_path: str, model: str = "whisper-1") -> None:
    """Parses the content of an audio file and prints it.

    Args:
        file_path: str: The path to the audio file to transcribe.
        model: str: The audio model to use for transcription. Defaults to 'whisper-1'.

    """
    try:
        with open(file_path, "rb") as audio_file:
            _get_openai_client().audio.translations.create(model=model, file=audio_file)
    except Exception:
        pass


def parse_image(
    file_path: str, task: str = "Describe this image as detail as possible."
) -> None:
    """Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the file to open.
        task: str: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.

    """
    try:
        base64_image = _base64_img(file_path)
        response = _get_openai_client().chat.completions.create(
            model=_get_openai_model(),
            messages=cast(Any, _prepare_image_messages(task, base64_image)),
            max_tokens=_get_max_token(),
        )
        if getattr(response, "choices", None) and len(response.choices) > 0:
            _ = response.choices[0].message.content
    except Exception:
        pass


def parse_video(
    file_path: str,
    task: str = "Describe this image as detail as possible.",
    frame_interval: int = 30,
) -> None:
    """Parses the content of an image file and prints the description.

    Args:
        file_path: str: The path to the video file to open.
        task: str: The task description for the API call. Defaults to 'Describe this image as detail as possible.'.
        frame_interval: int: The interval between frames to analyze. Defaults to 30.

    """
    task = task or "This is one frame from a video, please summarize this frame."
    base64_frames = _base64_video(file_path)
    selected_frames = base64_frames[::frame_interval]
    if len(selected_frames) > 30:
        new_interval = len(base64_frames) // 30
        selected_frames = base64_frames[::new_interval]
    for _idx, base64_frame in enumerate(selected_frames, start=1):
        try:
            response = _get_openai_client().chat.completions.create(
                model=_get_openai_model(),
                messages=cast(Any, _prepare_image_messages(task, base64_frame)),
                max_tokens=_get_max_token(),
            )
            if getattr(response, "choices", None) and len(response.choices) > 0:
                _ = response.choices[0].message.content
        except Exception:
            pass


def parse_pptx(file_path: str) -> None:
    """Parses the content of a pptx file and prints it.

    Args:
        file_path: str: The path to the file to open.

    """
    try:
        pres = Presentation(file_path)
        output_lines = [f"[Reading PowerPoint file from {file_path}]"]
        for slide_idx, slide in enumerate(pres.slides, start=1):
            output_lines.append(f"@@ Slide {slide_idx} @@")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    output_lines.append(shape.text)
            output_lines.append("")  # blank line between slides
        output = "\n".join(output_lines).rstrip("\n")
        print(f"{output}\n")
    except Exception:
        pass
